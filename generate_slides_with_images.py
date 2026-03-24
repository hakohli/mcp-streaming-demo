"""Generate PPTX slides for MCP Live presentation with embedded diagram images."""

import os, io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ──────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT    = RGBColor(0x00, 0xD4, 0xAA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY= RGBColor(0xBB, 0xBB, 0xCC)

BG_HEX    = "#1B1B2F"
ACCENT_HEX= "#00D4AA"
WHITE_HEX = "#FFFFFF"
GRAY_HEX  = "#BBBBCC"
ORANGE_HEX= "#FF8C00"
BOX_HEX   = "#2A2A4A"
RED_HEX   = "#FF4444"
BLUE_HEX  = "#4488FF"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helpers ──────────────────────────────────────────────────────────────
def set_bg(slide):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = BG_DARK

def add_text(slide, text, left, top, width, height, size=28, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.font.name = "Calibri"; p.alignment = align
    return tf

def add_bullets(slide, items, left, top, width, height, size=18, color=WHITE):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = "Calibri"; p.space_after = Pt(8)

def fig_to_stream(fig):
    buf = io.BytesIO(); fig.savefig(buf, format="png", dpi=180, bbox_inches="tight", facecolor=BG_HEX)
    plt.close(fig); buf.seek(0); return buf

def add_image_from_fig(slide, fig, left, top, width, height):
    slide.shapes.add_picture(fig_to_stream(fig), left, top, width, height)

def _box(ax, x, y, w, h, label, color=ACCENT_HEX, fs=11):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.15",
                                facecolor=BOX_HEX, edgecolor=color, linewidth=2))
    ax.text(x + w/2, y + h/2, label, ha="center", va="center",
            color=WHITE_HEX, fontsize=fs, fontweight="bold")

def _arrow(ax, x1, y1, x2, y2, color=ACCENT_HEX):
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle="-|>", color=color, lw=2))

# ── Diagram generators ──────────────────────────────────────────────────

def make_full_architecture_diagram():
    """Detailed component architecture matching the Mermaid diagram in README."""
    fig, ax = plt.subplots(figsize=(10, 5.5), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 14); ax.set_ylim(0, 8); ax.axis("off")

    # ── Event Sources (left) ──
    ax.text(1.5, 7.4, "Event Sources", ha="center", color=GRAY_HEX, fontsize=10, fontstyle="italic")
    _box(ax, 0.3, 6, 2.4, 1, "Log Simulator\nlog_simulator.py", GRAY_HEX, 9)

    # ── Kafka (center-left) ──
    ax.text(5, 7.4, "Message Bus", ha="center", color=ORANGE_HEX, fontsize=10, fontstyle="italic")
    _box(ax, 3.5, 6, 3, 1, "Apache Kafka\nTopic: app-logs\nConfluent 7.6.0", ORANGE_HEX, 9)

    # ── Streaming MCP Server (center) ──
    ax.text(5, 5, "Streaming MCP Server", ha="center", color=ACCENT_HEX, fontsize=11, fontweight="bold")
    # Server sub-components
    _box(ax, 0.3, 3.2, 2.2, 0.9, "Kafka\nConsumer", ORANGE_HEX, 9)
    _box(ax, 3.1, 3.2, 2.2, 0.9, "Anomaly\nDetection", BLUE_HEX, 9)
    _box(ax, 5.9, 3.2, 2.2, 0.9, "WebSocket\nServer :8765", ACCENT_HEX, 9)
    _box(ax, 8.7, 3.2, 2.2, 0.9, "Cedar Policy\nEngine", RED_HEX, 9)
    # Server internal arrows
    _arrow(ax, 2.6, 3.65, 3.0, 3.65, GRAY_HEX)
    _arrow(ax, 5.4, 3.65, 5.8, 3.65, GRAY_HEX)
    _arrow(ax, 8.2, 3.65, 8.6, 3.65, GRAY_HEX)
    # Bounding box for server
    ax.add_patch(FancyBboxPatch((0.1, 2.9), 11, 2.3, boxstyle="round,pad=0.2",
                                facecolor="none", edgecolor=ACCENT_HEX, linewidth=1.5, linestyle="--"))

    # ── AI Agent (right) ──
    ax.text(12.5, 7.4, "AI Agent", ha="center", color=ACCENT_HEX, fontsize=10, fontstyle="italic")
    _box(ax, 11.2, 6, 2.6, 1, "Subscribe\nto stream", ACCENT_HEX, 9)
    _box(ax, 11.2, 1.5, 2.6, 0.9, "Remediation\nError -> Fix", ACCENT_HEX, 9)
    _box(ax, 11.2, 0.2, 2.6, 0.9, "Sliding Window\n20-event monitor", BLUE_HEX, 9)

    # ── Connecting arrows ──
    # Simulator -> Kafka
    _arrow(ax, 2.8, 6.5, 3.4, 6.5, ORANGE_HEX)
    ax.text(3.1, 6.9, "produce", ha="center", color=ORANGE_HEX, fontsize=8)
    # Kafka -> Consumer
    _arrow(ax, 5, 6.0, 1.4, 4.2, ORANGE_HEX)
    ax.text(2.5, 5.3, "consume", ha="center", color=ORANGE_HEX, fontsize=8)
    # WebSocket -> Agent (push)
    _arrow(ax, 8.2, 3.9, 11.1, 6.3, ACCENT_HEX)
    ax.text(10, 5.5, "push: live\nevents", ha="center", color=ACCENT_HEX, fontsize=8)
    # Agent -> WebSocket (pull)
    _arrow(ax, 11.1, 1.9, 7.5, 3.2, BLUE_HEX)
    ax.text(9, 2.2, "request:\nget_anomalies()", ha="center", color=BLUE_HEX, fontsize=8)
    # Agent internal
    _arrow(ax, 12.5, 6.0, 12.5, 2.5, GRAY_HEX)

    # ── Legend ──
    ax.text(0.3, 0.3, "push (streaming)", color=ACCENT_HEX, fontsize=9)
    ax.text(3.3, 0.3, "pull (request/response)", color=BLUE_HEX, fontsize=9)
    ax.text(6.8, 0.3, "Kafka produce/consume", color=ORANGE_HEX, fontsize=9)

    return fig

def make_stale_context_diagram():
    fig, ax = plt.subplots(figsize=(8, 3.5), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 10); ax.set_ylim(0, 4); ax.axis("off")
    # Timeline
    ax.plot([1, 9], [2, 2], color=GRAY_HEX, lw=2)
    marks = [(2, "2:03pm\nDeploy"), (4, "2:04pm\nErrors spike"), (7, "2:15pm\nAgent polls")]
    for mx, label in marks:
        ax.plot(mx, 2, "o", color=ORANGE_HEX if "Error" in label else ACCENT_HEX, ms=10)
        ax.text(mx, 1.2, label, ha="center", va="top", color=WHITE_HEX, fontsize=10)
    # Blind spot
    ax.annotate("", xy=(7, 2.7), xytext=(4, 2.7),
                arrowprops=dict(arrowstyle="<->", color=RED_HEX, lw=2))
    ax.text(5.5, 3.1, "11 min blind spot", ha="center", color=RED_HEX, fontsize=12, fontweight="bold")
    return fig

def make_architecture_diagram():
    fig, ax = plt.subplots(figsize=(9, 3.5), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")
    # Pull model (top)
    ax.text(6, 4.6, "Traditional MCP (Pull)", ha="center", color=GRAY_HEX, fontsize=11)
    _box(ax, 1, 3.5, 2, 0.8, "Client"); _box(ax, 5, 3.5, 2, 0.8, "Server"); _box(ax, 9, 3.5, 2, 0.8, "Response")
    _arrow(ax, 3.1, 3.9, 4.9, 3.9, GRAY_HEX); _arrow(ax, 7.1, 3.9, 8.9, 3.9, GRAY_HEX)
    # Push model (bottom)
    ax.text(6, 2.6, "Streaming MCP (Push)", ha="center", color=ACCENT_HEX, fontsize=11, fontweight="bold")
    _box(ax, 0.3, 1, 2, 0.8, "Events"); _box(ax, 3.3, 1, 2, 0.8, "Kafka"); _box(ax, 6.3, 1, 2.5, 0.8, "MCP Server")
    _box(ax, 9.8, 1, 2, 0.8, "AI Agent")
    _arrow(ax, 2.4, 1.4, 3.2, 1.4); _arrow(ax, 5.4, 1.4, 6.2, 1.4); _arrow(ax, 8.9, 1.4, 9.7, 1.4)
    ax.text(7.55, 0.4, "WebSocket", ha="center", color=ACCENT_HEX, fontsize=9)
    return fig

def make_subscribe_protocol_diagram():
    fig, ax = plt.subplots(figsize=(7, 4), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis("off")
    # Agent and Server columns
    ax.plot([3, 3], [0.5, 5.5], color=ACCENT_HEX, lw=2, ls="--")
    ax.plot([7, 7], [0.5, 5.5], color=ACCENT_HEX, lw=2, ls="--")
    ax.text(3, 5.7, "Agent", ha="center", color=ACCENT_HEX, fontsize=13, fontweight="bold")
    ax.text(7, 5.7, "MCP Server", ha="center", color=ACCENT_HEX, fontsize=13, fontweight="bold")
    msgs = [
        (5.0, "subscribe(app-logs)", "->", ACCENT_HEX),
        (4.2, "event: log entry", "<-", ORANGE_HEX),
        (3.4, "event: error spike!", "<-", RED_HEX),
        (2.6, "get_anomalies()", "->", BLUE_HEX),
        (1.8, "anomaly summary", "<-", BLUE_HEX),
        (1.0, "unsubscribe", "->", GRAY_HEX),
    ]
    for y, label, direction, color in msgs:
        if direction == "->":
            _arrow(ax, 3.2, y, 6.8, y, color)
            ax.text(5, y + 0.2, label, ha="center", color=color, fontsize=10)
        else:
            _arrow(ax, 6.8, y, 3.2, y, color)
            ax.text(5, y + 0.2, label, ha="center", color=color, fontsize=10)
    return fig

def make_cedar_diagram():
    fig, ax = plt.subplots(figsize=(8, 3.5), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")
    _box(ax, 0.5, 2, 2.2, 1, "Agent\nRequest", BLUE_HEX)
    _box(ax, 4, 2, 2.5, 1, "Cedar\nPolicy Engine", ORANGE_HEX)
    _box(ax, 8, 3, 2.5, 0.8, "✓ ALLOW", ACCENT_HEX)
    _box(ax, 8, 1.2, 2.5, 0.8, "✗ DENY", RED_HEX)
    _arrow(ax, 2.8, 2.5, 3.9, 2.5, WHITE_HEX)
    _arrow(ax, 6.6, 2.8, 7.9, 3.4, ACCENT_HEX)
    _arrow(ax, 6.6, 2.2, 7.9, 1.6, RED_HEX)
    ax.text(6, 0.4, "Default Deny — no matching policy = blocked", ha="center",
            color=RED_HEX, fontsize=11, fontstyle="italic")
    return fig

def make_demo_flow_diagram():
    fig, ax = plt.subplots(figsize=(9, 3.5), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 12); ax.set_ylim(0, 4); ax.axis("off")
    boxes = [("Log\nSimulator", 0.3), ("Kafka", 2.8), ("MCP\nServer", 5.3), ("AI\nAgent", 7.8), ("Fix\nSuggestion", 10.3)]
    for label, x in boxes:
        c = RED_HEX if "Fix" in label else ACCENT_HEX
        _box(ax, x, 1.5, 1.8, 1, label, c)
    for i in range(len(boxes)-1):
        _arrow(ax, boxes[i][1]+1.9, 2, boxes[i+1][1]-0.1, 2)
    ax.text(6, 0.6, "< 2 seconds end-to-end", ha="center", color=ACCENT_HEX, fontsize=12, fontweight="bold")
    return fig

def make_anomaly_detection_diagram():
    fig, ax = plt.subplots(figsize=(8, 3.5), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 10); ax.set_ylim(0, 4); ax.axis("off")
    # Simulated error rate bar chart
    import numpy as np
    x = np.arange(10)
    rates = [5, 8, 6, 10, 12, 35, 55, 70, 45, 30]
    colors = [ACCENT_HEX if r < 30 else RED_HEX for r in rates]
    ax_inner = fig.add_axes([0.1, 0.15, 0.8, 0.7])
    ax_inner.set_facecolor(BOX_HEX)
    ax_inner.bar(x, rates, color=colors, edgecolor=BG_HEX, width=0.7)
    ax_inner.axhline(y=30, color=ORANGE_HEX, ls="--", lw=2, label="30% threshold")
    ax_inner.set_ylabel("Error Rate %", color=WHITE_HEX, fontsize=10)
    ax_inner.set_xlabel("Time Window", color=WHITE_HEX, fontsize=10)
    ax_inner.set_title("Anomaly Detection — Sliding Window", color=ACCENT_HEX, fontsize=12, fontweight="bold")
    ax_inner.tick_params(colors=GRAY_HEX)
    ax_inner.legend(facecolor=BOX_HEX, edgecolor=GRAY_HEX, labelcolor=WHITE_HEX)
    for spine in ax_inner.spines.values(): spine.set_color(GRAY_HEX)
    ax.axis("off")
    return fig

def make_use_cases_diagram():
    fig, ax = plt.subplots(figsize=(8, 4), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis("off")
    # Central hub
    _box(ax, 3.5, 2.5, 3, 1, "Streaming\nMCP", ACCENT_HEX, 13)
    cases = [
        (0.2, 4.5, "DevOps", ACCENT_HEX),
        (4, 4.8, "Trading", ORANGE_HEX),
        (7.5, 4.5, "Security", RED_HEX),
        (0.2, 0.5, "CI/CD", BLUE_HEX),
        (7.5, 0.5, "IoT", ORANGE_HEX),
    ]
    for x, y, label, color in cases:
        _box(ax, x, y, 2.2, 0.8, label, color, 10)
        cx, cy = x + 1.1, y + 0.4
        _arrow(ax, 5, 3 if y > 3 else 2.5, cx, cy + (-.4 if y > 3 else .4), color)
    return fig

def make_takeaways_diagram():
    fig, ax = plt.subplots(figsize=(8, 3.5), facecolor=BG_HEX)
    ax.set_facecolor(BG_HEX); ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis("off")
    items = [
        (1, 3.8, "Kafka", "Durable\nStreaming"),
        (4, 3.8, "WebSocket", "Real-time\nPush"),
        (7, 3.8, "Cedar", "AuthZ\nPolicies"),
        (4, 1, "Streaming\nMCP Server", ""),
    ]
    colors = [ORANGE_HEX, BLUE_HEX, RED_HEX, ACCENT_HEX]
    for (x, y, label, sub), c in zip(items, colors):
        _box(ax, x, y, 2, 1, f"{label}\n{sub}" if sub else label, c, 10)
    _arrow(ax, 2, 3.8, 4.2, 2.1, ORANGE_HEX)
    _arrow(ax, 5, 3.8, 5, 2.1, BLUE_HEX)
    _arrow(ax, 8, 3.8, 5.8, 2.1, RED_HEX)
    ax.text(5, 0.3, "Hybrid Push + Pull on a Single Connection", ha="center",
            color=ACCENT_HEX, fontsize=12, fontweight="bold")
    return fig

# ── Slide definitions ────────────────────────────────────────────────────
SLIDES = [
    {"title": "MCP Live:\nStreaming Context to AI Agents",
     "subtitle": "Building Reactive MCP Servers with Kafka, WebSockets & Cedar\ngithub.com/hakohli/mcp-streaming-demo", "type": "title",
     "notes": (
         "TALKING POINTS:\n"
         "- Welcome everyone. Today we're going to talk about extending the Model Context Protocol — MCP — beyond its traditional request/response model.\n"
         "- The core idea: what if MCP servers could push real-time context to AI agents instead of waiting to be asked?\n"
         "- We'll walk through a working implementation using three open-source technologies: Apache Kafka for durable event streaming, WebSockets for real-time delivery, and Cedar for fine-grained authorization.\n"
         "- By the end, you'll see a live demo of an AI agent detecting and responding to a production incident in under 2 seconds — something impossible with traditional polling.\n"
         "- Everything we show today is open source and available in the repo."
     )},
    {"title": "System Architecture",
     "bullets": [
         "Log Simulator produces events to Kafka (app-logs topic)",
         "MCP Server consumes from Kafka, maintains rolling window",
         "Anomaly Detection tracks error rates server-side",
         "WebSocket Server pushes live events to subscribed agents",
         "Cedar Policy Engine enforces permit/forbid on every action",
         "AI Agent subscribes, monitors 20-event sliding window",
         "Remediation Engine maps error patterns to fix suggestions",
         "Hybrid: push (streaming) + pull (get_anomalies) on same connection",
     ], "diagram": make_full_architecture_diagram, "font_size": 16,
     "notes": (
         "TALKING POINTS:\n"
         "- This is the full system architecture. Let me walk through each component.\n"
         "- Starting on the left: the Log Simulator generates realistic application logs — INFO, WARN, ERROR, FATAL — from five simulated services: api-gateway, auth-service, payment-service, order-service, and inventory-db. Every ~30 seconds it triggers an error spike.\n"
         "- Events flow into Apache Kafka (Confluent 7.6.0, running in KRaft mode via Docker). Kafka gives us durability, replayability, and decoupling between producers and consumers.\n"
         "- The Streaming MCP Server is the core — it has four sub-components:\n"
         "  1. Kafka Consumer: pulls events and maintains a rolling 100-event window for analysis.\n"
         "  2. Anomaly Detection: tracks error counts by pattern (ConnectionRefused, OOM, Timeout, etc.) across the window.\n"
         "  3. WebSocket Server on port 8765: pushes live events to all subscribed agents. Each subscriber gets a bounded 500-event queue — if they fall behind, oldest events are dropped.\n"
         "  4. Cedar Policy Engine: evaluates every agent request against declarative permit/forbid policies. Default deny — no matching policy means the request is blocked.\n"
         "- On the right: the AI Agent. It subscribes once and receives a continuous stream. It runs a 20-event sliding window — when error rate exceeds 30%, it flags an anomaly. It can also send get_anomalies() back to the server for a structured summary — that's the hybrid protocol, push AND pull on the same WebSocket.\n"
         "- The Remediation Engine maps known error patterns to actionable fixes: ConnectionRefused → check service status; OOM → increase heap; Timeout → circuit breaker; Deadlock → review isolation levels.\n"
         "- Notice the two arrow colors: teal for streaming push, blue for request/response pull. Both coexist on the same connection. This is the key design insight."
     )},
    {"title": "The Stale Context Problem",
     "bullets": [
         "Traditional MCP = request/response snapshots",
         "Agent asks for logs → gets point-in-time dump → works with stale data",
         "Deploy at 2:03pm, errors spike at 2:04pm, agent doesn't know until 2:15pm",
         "What if context came TO the agent, as it happened?",
     ], "diagram": make_stale_context_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- Let's start with the problem. Today, MCP works like a snapshot camera — the agent asks a question, gets a point-in-time answer, and moves on.\n"
         "- Imagine this real scenario: you deploy at 2:03pm. Errors start spiking at 2:04pm. But your agent doesn't poll again until 2:15pm. That's an 11-minute blind spot where your production system is on fire and nobody knows.\n"
         "- The diagram on the right shows this timeline visually — the red gap is the danger zone.\n"
         "- In fast-moving environments — incident response, trading, security — 11 minutes is an eternity.\n"
         "- The fundamental question: can we flip the model so context flows TO the agent continuously, rather than the agent having to ask for it?\n"
         "- Spoiler: yes, and it's not as hard as you'd think."
     )},
    {"title": "Streaming MCP Architecture",
     "bullets": [
         "Standard MCP: Client → Server → Response (pull model)",
         "Streaming MCP: Server → Client via subscriptions (push model)",
         "Why Kafka: durable, replayable, backpressure handling",
         "Hybrid protocol: streaming + request/response on same connection",
     ], "diagram": make_architecture_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- Here's the architecture side by side. Top row is traditional MCP — pull-based, request/response.\n"
         "- Bottom row is what we're building: event sources feed into Kafka, Kafka feeds the MCP server, and the server pushes to agents over WebSockets.\n"
         "- Why Kafka specifically? Three reasons: (1) durability — events are persisted, so if an agent disconnects and reconnects, it can replay what it missed. (2) Backpressure — Kafka handles slow consumers gracefully via consumer groups and offsets. (3) It's probably already in your infrastructure.\n"
         "- The key insight is the hybrid protocol — we keep standard request/response for on-demand queries, but ADD streaming on the same WebSocket connection. The agent can subscribe to a live feed AND still ask ad-hoc questions.\n"
         "- This means zero changes to existing MCP tools — we're extending, not replacing."
     )},
    {"title": "Key Protocol Extension: subscribe",
     "bullets": [
         "New method: subscribe / unsubscribe for continuous push",
         "Still supports get_anomalies, get_context (standard request/response)",
         "Bounded queues (500 events) handle backpressure automatically",
     ], "diagram": make_subscribe_protocol_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- The protocol extension is minimal — we add just two new methods: subscribe and unsubscribe.\n"
         "- Walk through the sequence diagram: Agent sends subscribe('app-logs'). Server starts pushing events as they arrive from Kafka. Normal log entries flow through. Then an error spike hits — the agent sees it immediately.\n"
         "- Here's the hybrid part: mid-stream, the agent sends get_anomalies() — a standard request/response call — and gets back a structured anomaly summary. Both protocols coexist on the same WebSocket.\n"
         "- When the agent is done, it sends unsubscribe to clean up.\n"
         "- Backpressure is handled with bounded queues — each subscriber gets a 500-event buffer. If a subscriber falls behind, oldest events are dropped rather than letting memory grow unbounded.\n"
         "- Dead subscribers are auto-cleaned — if a WebSocket disconnects, the server detects it and removes the subscription. No resource leaks."
     )},
    {"title": "Cedar Authorization for MCP",
     "bullets": [
         "Open-source policy language by AWS",
         "Declarative: permit/forbid(principal, action, resource)",
         "Role-based access control for MCP actions",
         "Default deny — no matching policy = request denied",
     ], "diagram": make_cedar_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- Now, streaming data to agents raises an obvious question: who's allowed to see what?\n"
         "- Cedar is an open-source authorization policy language created by AWS. It's now being adopted in the MCP ecosystem — used in Bedrock AgentCore and ToolHive for MCP tool access control.\n"
         "- The model is simple: every request is evaluated as (principal, action, resource). You write permit and forbid policies declaratively.\n"
         "- For our MCP server, actions include: subscribe, get_anomalies, get_context, call_tool. Resources are streams like 'app-logs' or 'audit-logs'.\n"
         "- Critical design choice: default deny. If no policy matches a request, it's blocked. This is the secure-by-default posture you want for production.\n"
         "- The diagram shows the flow: request comes in, Cedar evaluates policies, result is either ALLOW or DENY. No ambiguity."
     )},
    {"title": "Cedar Policy Examples",
     "bullets": [
         'permit(principal, action == Action::"subscribe", resource == Stream::"app-logs");',
         'permit(principal, action == Action::"get_anomalies", resource)\n    when { principal.role == "ops-agent" };',
         'forbid(principal, action == Action::"subscribe", resource == Stream::"audit-logs")\n    when { principal.role == "readonly" };',
         "Forbid always wins over permit — secure by default",
     ], "font_size": 16,
     "notes": (
         "TALKING POINTS:\n"
         "- Let's look at real Cedar policies. These are actual policies from our demo.\n"
         "- First policy: anyone can subscribe to the app-logs stream. This is your general observability data.\n"
         "- Second policy: only agents with the 'ops-agent' role can call get_anomalies. You don't want every agent querying for anomaly summaries — that's an expensive operation.\n"
         "- Third policy: explicitly forbid readonly agents from subscribing to audit-logs. This is sensitive compliance data.\n"
         "- Key rule: forbid ALWAYS wins over permit. Even if another policy permits access, a single forbid blocks it. This prevents accidental over-permissioning.\n"
         "- These policies are human-readable, auditable, and version-controllable. Your security team can review them in a PR."
     )},
    {"title": "Live Demo: Real-time Log Monitoring",
     "bullets": [
         "Log simulator → Kafka → MCP Server → AI Agent",
         "Error spike → agent detects anomaly within 1 second",
         "Agent queries for anomaly summary mid-stream",
         "Cedar blocks unauthorized agents from sensitive streams",
     ], "diagram": make_demo_flow_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- Time for the live demo. Stack: Confluent Kafka 7.6.0 in Docker, Python 3.9, confluent-kafka 2.13.2, websockets 15.0.1.\n"
         "- Here's what's running: a log simulator generates realistic application logs and pushes them to Kafka. The streaming MCP server consumes from Kafka and pushes to connected agents.\n"
         "- We'll see normal traffic flowing first — INFO and DEBUG level logs streaming through in real-time.\n"
         "- Then we'll trigger an error spike — simulating a production incident like a database connection failure.\n"
         "- Watch the agent: it detected '🚨 ANOMALY DETECTED — error rate 60% in last 20 events' — within 1 second of the spike.\n"
         "- The agent automatically suggested fixes: ConnectionRefused → check service status; OutOfMemoryError → increase JVM heap; Timeout → add circuit breaker; Deadlock → review transaction isolation.\n"
         "- We'll also show Cedar in action: try subscribing with an unauthorized agent and watch it get denied.\n"
         "- End-to-end latency from error occurrence to fix suggestion: under 2 seconds."
     )},
    {"title": "Anomaly Detection Flow",
     "bullets": [
         "Rolling window of last 100 events maintained server-side",
         "Error rate > 30% triggers anomaly alert",
         "Agent maps error patterns to fix suggestions automatically",
     ], "diagram": make_anomaly_detection_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- Let's break down how anomaly detection works in this system.\n"
         "- Server-side: we maintain a rolling window of the last 100 events. This gives us a stable baseline without unbounded memory growth.\n"
         "- Agent-side: the agent tracks errors in a sliding 20-event window. When the error rate exceeds 30%, it triggers an anomaly alert.\n"
         "- The bar chart on the right shows this visually — green bars are normal error rates, red bars are above the 30% threshold. The orange dashed line is the trigger point.\n"
         "- When triggered, the agent calls get_anomalies() on the MCP server, which returns the top-N error patterns with counts and timestamps.\n"
         "- The agent then maps these patterns to remediation suggestions: ConnectionError → check database connectivity and restart. OutOfMemoryError → increase heap allocation. SSLError → renew certificates.\n"
         "- All of this happens within 1-2 seconds of the actual error. Compare that to the 11-minute blind spot we started with."
     )},
    {"title": "Production Patterns & Gotchas",
     "bullets": [
         "1. Backpressure: bounded queues, drop slow subscribers",
         "2. Hybrid protocol: streaming AND on-demand on same WebSocket",
         "3. Kafka consumer groups: same = load balance, different = fan-out",
         "4. Reconnection: exponential backoff + context snapshot",
         "5. Be selective: stream errors & anomalies, batch normal metrics",
     ],
     "notes": (
         "TALKING POINTS:\n"
         "- These are the lessons learned from building this. Things that will bite you in production.\n"
         "- Backpressure is #1 for a reason. If an agent can't keep up with the event stream, you need bounded queues. Our implementation uses a 500-event asyncio queue per subscriber. When full, oldest events are dropped. You could also use dead-letter topics in Kafka for events that couldn't be delivered.\n"
         "- The hybrid protocol is the sweet spot. Don't force everything into streaming OR request/response. Let agents subscribe to live feeds for continuous monitoring, but also query on-demand for summaries and context. Same WebSocket, same connection.\n"
         "- Kafka consumer groups are powerful but confusing. Same group ID = events are load-balanced across consumers (each event goes to ONE consumer). Different group IDs = fan-out (each event goes to ALL consumers). Choose based on your use case.\n"
         "- Reconnection strategy matters. Use exponential backoff to avoid thundering herd. On reconnect, send a context snapshot so the agent doesn't start from zero.\n"
         "- Be selective about what you stream. Errors and anomalies? Stream them. Normal INFO-level metrics? Batch those and let agents poll. Don't flood the pipe with noise."
     )},
    {"title": "What Streaming MCP Unlocks",
     "bullets": [
         "DevOps — real-time incident detection and auto-remediation",
         "Trading — live market data feeding AI decision agents",
         "Security — streaming audit logs with Cedar-authorized access",
         "CI/CD — file watchers notifying agents of code changes",
         "IoT — sensor data streaming to predictive maintenance agents",
     ], "diagram": make_use_cases_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- This pattern isn't just for log monitoring. Once you have streaming MCP, a whole category of use cases opens up.\n"
         "- DevOps is the obvious one — what we demoed today. Real-time incident detection, automated triage, and remediation suggestions. Imagine on-call agents that never sleep.\n"
         "- Trading: feed live market data — price ticks, order book changes, news events — directly to AI decision agents. Latency matters here; polling is not an option.\n"
         "- Security: stream audit logs to compliance agents with Cedar controlling who sees what. Your SOC agent gets real-time alerts; your readonly dashboard agent gets summaries only.\n"
         "- CI/CD: file watchers that notify agents when code changes. Agent sees a PR merged, automatically runs tests, checks for security issues, updates documentation.\n"
         "- IoT: sensor data from factory floors, vehicles, or infrastructure streaming to predictive maintenance agents. Detect anomalies before equipment fails.\n"
         "- The common thread: any domain where waiting to be asked is too slow."
     )},
    {"title": "Key Takeaways",
     "bullets": [
         "MCP doesn't have to be request/response only",
         "Kafka (Confluent 7.6.0) + WebSocket = durable, scalable streaming bridge",
         "Cedar provides declarative, auditable authorization for MCP",
         "Hybrid protocol (push + pull) on single connection is the sweet spot",
         "All open source — Python 3.9+, Docker, ~200 lines of code",
     ], "diagram": make_takeaways_diagram,
     "notes": (
         "TALKING POINTS:\n"
         "- Four things to remember from this talk.\n"
         "- First: MCP is extensible. The spec doesn't limit you to request/response. Streaming is a natural extension that the protocol can support.\n"
         "- Second: you don't need exotic infrastructure. Kafka gives you durable, replayable event streaming. WebSockets give you real-time push. Both are battle-tested, widely deployed, and well-understood.\n"
         "- Third: authorization is not optional. Cedar gives you declarative, auditable policies that your security team can review. Default deny means you start secure and open up access explicitly.\n"
         "- Fourth: the hybrid protocol is the key design decision. Don't choose between streaming and request/response — do both on the same connection. Agents get the best of both worlds.\n"
         "- The diagram shows how these three pieces converge into the streaming MCP server. Each solves a different problem; together they create something powerful."
     )},
    {"title": "Thank You & Q&A",
     "subtitle": "github.com/hakohli/mcp-streaming-demo\nAll open source: Kafka · Cedar · WebSockets · Python", "type": "title",
     "notes": (
         "TALKING POINTS:\n"
         "- Thank you for your time. The complete working code is on GitHub: github.com/hakohli/mcp-streaming-demo\n"
         "- Everything is open source: Confluent Kafka 7.6.0, Cedar, WebSockets 15.0.1, Python 3.9+.\n"
         "- To run the demo yourself: docker compose up -d for Kafka, then run streaming_mcp_server.py, log_simulator.py, and agent_client.py.\n"
         "- Happy to take questions on the architecture, Cedar policies, Kafka integration, or anything else.\n"
         "- If you're interested in contributing or extending this — for example, adding SSE transport, multi-topic subscriptions, or persistent agent state — let's talk."
     )},
]

# ── Build slides ─────────────────────────────────────────────────────────
for s in SLIDES:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Add speaker notes
    if s.get("notes"):
        slide.notes_slide.notes_text_frame.text = s["notes"]

    if s.get("type") == "title":
        add_text(slide, s["title"], Inches(1), Inches(2), Inches(11), Inches(2.5),
                 size=44, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        if s.get("subtitle"):
            add_text(slide, s["subtitle"], Inches(1), Inches(4.5), Inches(11), Inches(1.5),
                     size=24, color=LIGHT_GRAY, bold=False, align=PP_ALIGN.CENTER)
    else:
        add_text(slide, s["title"], Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                 size=32, color=ACCENT, bold=True)
        has_diagram = "diagram" in s
        bullet_w = Inches(5.5) if has_diagram else Inches(11)
        fs = s.get("font_size", 18)
        add_bullets(slide, s["bullets"], Inches(1.0), Inches(1.4), bullet_w, Inches(5.2), size=fs)
        if has_diagram:
            fig = s["diagram"]()
            add_image_from_fig(slide, fig, Inches(6.8), Inches(1.4), Inches(6), Inches(5))

out = "/Users/hakohli/mcp-streaming-demo/MCP_Live_Streaming_Context.pptx"
prs.save(out)
print(f"✅ Saved with images: {out}")
