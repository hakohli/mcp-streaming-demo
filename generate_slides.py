"""Generate PPTX slides for MCP Live presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_title_text(slide, text, left, top, width, height, size=28, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf

def add_bullets(slide, items, left, top, width, height, size=18, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return tf

SLIDES = [
    {
        "title": "MCP Live:\nStreaming Context to AI Agents",
        "subtitle": "Building Reactive MCP Servers with Kafka, WebSockets & Cedar",
        "type": "title",
    },
    {
        "title": "The Stale Context Problem",
        "bullets": [
            "Traditional MCP = request/response snapshots",
            "Agent asks for logs → gets point-in-time dump → works with stale data",
            "Deploy at 2:03pm, errors spike at 2:04pm, agent doesn't know until 2:15pm",
            "11 minutes of blind spot = missed production incidents",
            "What if context came TO the agent, as it happened?",
        ],
    },
    {
        "title": "Streaming MCP Architecture",
        "bullets": [
            "Standard MCP: Client → Server → Response (pull model)",
            "Streaming MCP: Server → Client via subscriptions (push model)",
            "Event Sources → Kafka → Streaming MCP Server → AI Agent (WebSocket)",
            "Why Kafka: durable, replayable, backpressure handling, already in your infra",
            "Hybrid protocol: streaming + request/response on same connection",
        ],
    },
    {
        "title": "Key Protocol Extension: subscribe",
        "bullets": [
            "New method: subscribe / unsubscribe for continuous push",
            "Agent subscribes once, receives events as they happen",
            "Still supports get_anomalies, get_context (standard request/response)",
            "Bounded queues (500 events) handle backpressure automatically",
            "Dead subscribers auto-cleaned — no resource leaks",
        ],
    },
    {
        "title": "Cedar Authorization for MCP",
        "bullets": [
            "Open-source policy language by AWS — now used in MCP authorization",
            "Declarative: permit/forbid(principal, action, resource) when { conditions }",
            "MCP actions: subscribe, get_anomalies, get_context, call_tool",
            "Role-based: ops-agent can subscribe, senior-ops gets full context",
            "Default deny — no matching policy = request denied",
            "Used in Bedrock AgentCore & ToolHive for MCP tool access control",
        ],
    },
    {
        "title": "Cedar Policy Examples",
        "bullets": [
            'permit(principal, action == Action::"subscribe", resource == Stream::"app-logs");',
            'permit(principal, action == Action::"get_anomalies", resource)\n    when { principal.role == "ops-agent" };',
            'forbid(principal, action == Action::"subscribe", resource == Stream::"audit-logs")\n    when { principal.role == "readonly" };',
            "Forbid always wins over permit — secure by default",
        ],
        "font_size": 16,
    },
    {
        "title": "Live Demo: Real-time Log Monitoring",
        "bullets": [
            "Log simulator → Kafka → Streaming MCP Server → AI Agent",
            "Normal traffic flows through in real-time",
            "Error spike hits → agent detects anomaly within 1 second",
            "Agent queries server for anomaly summary mid-stream",
            "Suggests fixes: restart service, increase heap, renew cert",
            "Cedar blocks unauthorized agents from sensitive streams",
        ],
    },
    {
        "title": "Anomaly Detection Flow",
        "bullets": [
            "Rolling window of last 100 events maintained server-side",
            "Agent tracks errors in sliding 20-event window",
            "Error rate > 30% triggers anomaly alert",
            "Server provides top-N error patterns with counts",
            "Agent maps error patterns to fix suggestions automatically",
            "All within 1–2 seconds of the actual error",
        ],
    },
    {
        "title": "Production Patterns & Gotchas",
        "bullets": [
            "1. Backpressure: bounded queues, drop slow subscribers, dead-letter topics",
            "2. Hybrid protocol: streaming AND on-demand on same WebSocket",
            "3. Kafka consumer groups: same group = load balance, different = fan-out",
            "4. Reconnection: exponential backoff + context snapshot on reconnect",
            "5. Be selective: stream errors & anomalies, batch normal metrics",
        ],
    },
    {
        "title": "What Streaming MCP Unlocks",
        "bullets": [
            "DevOps — real-time incident detection and auto-remediation",
            "Trading — live market data feeding AI decision agents",
            "Security — streaming audit logs with Cedar-authorized access",
            "CI/CD — file watchers notifying agents of code changes",
            "IoT — sensor data streaming to predictive maintenance agents",
        ],
    },
    {
        "title": "Key Takeaways",
        "bullets": [
            "MCP doesn't have to be request/response only",
            "Kafka + WebSocket = durable, scalable streaming bridge",
            "Cedar provides declarative, auditable authorization for MCP",
            "Hybrid protocol (push + pull) on single connection is the sweet spot",
            "All open source: Kafka, Cedar, WebSockets, Python",
        ],
    },
    {
        "title": "Thank You & Q&A",
        "subtitle": "Code: ~/mcp-streaming-demo\nAll open source: Kafka · Cedar · WebSockets",
        "type": "title",
    },
]

for s in SLIDES:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide)

    if s.get("type") == "title":
        add_title_text(slide, s["title"],
                       Inches(1), Inches(2), Inches(11), Inches(2.5),
                       size=44, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        if s.get("subtitle"):
            add_title_text(slide, s["subtitle"],
                           Inches(1), Inches(4.5), Inches(11), Inches(1.5),
                           size=24, color=LIGHT_GRAY, bold=False, align=PP_ALIGN.CENTER)
    else:
        add_title_text(slide, s["title"],
                       Inches(0.8), Inches(0.4), Inches(11), Inches(1),
                       size=32, color=ACCENT, bold=True)
        fs = s.get("font_size", 18)
        add_bullets(slide, s["bullets"],
                    Inches(1.0), Inches(1.6), Inches(11), Inches(5.2),
                    size=fs, color=WHITE)

out = "/Users/hakohli/mcp-streaming-demo/MCP_Live_Streaming_Context.pptx"
prs.save(out)
print(f"Saved: {out}")
