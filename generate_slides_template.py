"""Generate PPTX slides using MCP-DevSummit-PPT.pptx template."""

import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Template path ────────────────────────────────────────────────────────
TEMPLATE = "/Users/hakohli/Downloads/MCP-DevSummit-PPT.pptx"

# ── Diagram colours ──────────────────────────────────────────────────────
ACCENT_HEX = "#00D4AA"
WHITE_HEX  = "#FFFFFF"
GRAY_HEX   = "#BBBBCC"
ORANGE_HEX = "#FF8C00"
BOX_HEX    = "#2A2A4A"
RED_HEX    = "#FF4444"
BLUE_HEX   = "#4488FF"
TRANSPARENT = (0, 0, 0, 0)  # for fig background

def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=180, bbox_inches="tight", transparent=True)
    plt.close(fig); buf.seek(0); return buf

def _box(ax, x, y, w, h, label, color=ACCENT_HEX, fs=11):
    ax.add_patch(FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.15",
                                facecolor=BOX_HEX, edgecolor=color, linewidth=2))
    ax.text(x + w/2, y + h/2, label, ha="center", va="center",
            color=WHITE_HEX, fontsize=fs, fontweight="bold")

def _arrow(ax, x1, y1, x2, y2, color=ACCENT_HEX):
    ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                arrowprops=dict(arrowstyle="-|>", color=color, lw=2))

# ── Diagram generators ──────────────────────────────────────────────────

def make_full_architecture_diagram():
    fig, ax = plt.subplots(figsize=(9, 5), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 14); ax.set_ylim(0, 8); ax.axis("off")
    ax.text(1.5, 7.4, "Event Sources", ha="center", color=GRAY_HEX, fontsize=10, fontstyle="italic")
    _box(ax, 0.3, 6, 2.4, 1, "Log Simulator\nlog_simulator.py", GRAY_HEX, 9)
    ax.text(5, 7.4, "Message Bus", ha="center", color=ORANGE_HEX, fontsize=10, fontstyle="italic")
    _box(ax, 3.5, 6, 3, 1, "Apache Kafka\nTopic: app-logs\nConfluent 7.6.0", ORANGE_HEX, 9)
    ax.text(5, 5, "Streaming MCP Server", ha="center", color=ACCENT_HEX, fontsize=11, fontweight="bold")
    _box(ax, 0.3, 3.2, 2.2, 0.9, "Kafka\nConsumer", ORANGE_HEX, 9)
    _box(ax, 3.1, 3.2, 2.2, 0.9, "Anomaly\nDetection", BLUE_HEX, 9)
    _box(ax, 5.9, 3.2, 2.2, 0.9, "WebSocket\nServer :8765", ACCENT_HEX, 9)
    _box(ax, 8.7, 3.2, 2.2, 0.9, "Cedar Policy\nEngine", RED_HEX, 9)
    _arrow(ax, 2.6, 3.65, 3.0, 3.65, GRAY_HEX); _arrow(ax, 5.4, 3.65, 5.8, 3.65, GRAY_HEX)
    _arrow(ax, 8.2, 3.65, 8.6, 3.65, GRAY_HEX)
    ax.add_patch(FancyBboxPatch((0.1, 2.9), 11, 2.3, boxstyle="round,pad=0.2",
                                facecolor="none", edgecolor=ACCENT_HEX, linewidth=1.5, linestyle="--"))
    ax.text(12.5, 7.4, "AI Agent", ha="center", color=ACCENT_HEX, fontsize=10, fontstyle="italic")
    _box(ax, 11.2, 6, 2.6, 1, "Subscribe\nto stream", ACCENT_HEX, 9)
    _box(ax, 11.2, 1.5, 2.6, 0.9, "Remediation\nError -> Fix", ACCENT_HEX, 9)
    _box(ax, 11.2, 0.2, 2.6, 0.9, "Sliding Window\n20-event monitor", BLUE_HEX, 9)
    _arrow(ax, 2.8, 6.5, 3.4, 6.5, ORANGE_HEX); ax.text(3.1, 6.9, "produce", ha="center", color=ORANGE_HEX, fontsize=8)
    _arrow(ax, 5, 6.0, 1.4, 4.2, ORANGE_HEX); ax.text(2.5, 5.3, "consume", ha="center", color=ORANGE_HEX, fontsize=8)
    _arrow(ax, 8.2, 3.9, 11.1, 6.3, ACCENT_HEX); ax.text(10, 5.5, "push: live\nevents", ha="center", color=ACCENT_HEX, fontsize=8)
    _arrow(ax, 11.1, 1.9, 7.5, 3.2, BLUE_HEX); ax.text(9, 2.2, "request:\nget_anomalies()", ha="center", color=BLUE_HEX, fontsize=8)
    _arrow(ax, 12.5, 6.0, 12.5, 2.5, GRAY_HEX)
    ax.text(0.3, 0.3, "push (streaming)", color=ACCENT_HEX, fontsize=9)
    ax.text(3.3, 0.3, "pull (request/response)", color=BLUE_HEX, fontsize=9)
    ax.text(6.8, 0.3, "Kafka produce/consume", color=ORANGE_HEX, fontsize=9)
    return fig

def make_stale_context_diagram():
    fig, ax = plt.subplots(figsize=(8, 3), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 10); ax.set_ylim(0, 4); ax.axis("off")
    ax.plot([1, 9], [2, 2], color=GRAY_HEX, lw=2)
    for mx, label in [(2, "2:03pm\nDeploy"), (4, "2:04pm\nErrors spike"), (7, "2:15pm\nAgent polls")]:
        ax.plot(mx, 2, "o", color=ORANGE_HEX if "Error" in label else ACCENT_HEX, ms=10)
        ax.text(mx, 1.2, label, ha="center", va="top", color=WHITE_HEX, fontsize=10)
    ax.annotate("", xy=(7, 2.7), xytext=(4, 2.7), arrowprops=dict(arrowstyle="<->", color=RED_HEX, lw=2))
    ax.text(5.5, 3.1, "11 min blind spot", ha="center", color=RED_HEX, fontsize=12, fontweight="bold")
    return fig

def make_architecture_diagram():
    fig, ax = plt.subplots(figsize=(8, 3), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")
    ax.text(6, 4.6, "Traditional MCP (Pull)", ha="center", color=GRAY_HEX, fontsize=11)
    _box(ax, 1, 3.5, 2, 0.8, "Client"); _box(ax, 5, 3.5, 2, 0.8, "Server"); _box(ax, 9, 3.5, 2, 0.8, "Response")
    _arrow(ax, 3.1, 3.9, 4.9, 3.9, GRAY_HEX); _arrow(ax, 7.1, 3.9, 8.9, 3.9, GRAY_HEX)
    ax.text(6, 2.6, "Streaming MCP (Push)", ha="center", color=ACCENT_HEX, fontsize=11, fontweight="bold")
    _box(ax, 0.3, 1, 2, 0.8, "Events"); _box(ax, 3.3, 1, 2, 0.8, "Kafka"); _box(ax, 6.3, 1, 2.5, 0.8, "MCP Server")
    _box(ax, 9.8, 1, 2, 0.8, "AI Agent")
    _arrow(ax, 2.4, 1.4, 3.2, 1.4); _arrow(ax, 5.4, 1.4, 6.2, 1.4); _arrow(ax, 8.9, 1.4, 9.7, 1.4)
    ax.text(7.55, 0.4, "WebSocket", ha="center", color=ACCENT_HEX, fontsize=9)
    return fig

def make_subscribe_protocol_diagram():
    fig, ax = plt.subplots(figsize=(6, 3.5), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis("off")
    ax.plot([3, 3], [0.5, 5.5], color=ACCENT_HEX, lw=2, ls="--")
    ax.plot([7, 7], [0.5, 5.5], color=ACCENT_HEX, lw=2, ls="--")
    ax.text(3, 5.7, "Agent", ha="center", color=ACCENT_HEX, fontsize=13, fontweight="bold")
    ax.text(7, 5.7, "MCP Server", ha="center", color=ACCENT_HEX, fontsize=13, fontweight="bold")
    for y, label, d, c in [(5.0,"subscribe(app-logs)","->",ACCENT_HEX),(4.2,"event: log entry","<-",ORANGE_HEX),
                            (3.4,"event: error spike!","<-",RED_HEX),(2.6,"get_anomalies()","->",BLUE_HEX),
                            (1.8,"anomaly summary","<-",BLUE_HEX),(1.0,"unsubscribe","->",GRAY_HEX)]:
        if d == "->": _arrow(ax, 3.2, y, 6.8, y, c)
        else: _arrow(ax, 6.8, y, 3.2, y, c)
        ax.text(5, y + 0.2, label, ha="center", color=c, fontsize=10)
    return fig

def make_cedar_diagram():
    fig, ax = plt.subplots(figsize=(8, 3), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 12); ax.set_ylim(0, 5); ax.axis("off")
    _box(ax, 0.5, 2, 2.2, 1, "Agent\nRequest", BLUE_HEX)
    _box(ax, 4, 2, 2.5, 1, "Cedar\nPolicy Engine", ORANGE_HEX)
    _box(ax, 8, 3, 2.5, 0.8, "ALLOW", ACCENT_HEX); _box(ax, 8, 1.2, 2.5, 0.8, "DENY", RED_HEX)
    _arrow(ax, 2.8, 2.5, 3.9, 2.5, WHITE_HEX)
    _arrow(ax, 6.6, 2.8, 7.9, 3.4, ACCENT_HEX); _arrow(ax, 6.6, 2.2, 7.9, 1.6, RED_HEX)
    ax.text(6, 0.4, "Default Deny — no matching policy = blocked", ha="center", color=RED_HEX, fontsize=11, fontstyle="italic")
    return fig

def make_demo_flow_diagram():
    fig, ax = plt.subplots(figsize=(8, 3), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 12); ax.set_ylim(0, 4); ax.axis("off")
    for label, x in [("Log\nSimulator",0.3),("Kafka",2.8),("MCP\nServer",5.3),("AI\nAgent",7.8),("Fix\nSuggestion",10.3)]:
        _box(ax, x, 1.5, 1.8, 1, label, RED_HEX if "Fix" in label else ACCENT_HEX)
    for x1, x2 in [(0.3,2.8),(2.8,5.3),(5.3,7.8),(7.8,10.3)]:
        _arrow(ax, x1+1.9, 2, x2-0.1, 2)
    ax.text(6, 0.6, "< 2 seconds end-to-end", ha="center", color=ACCENT_HEX, fontsize=12, fontweight="bold")
    return fig

def make_anomaly_detection_diagram():
    import numpy as np
    fig = plt.figure(figsize=(7, 3), facecolor=TRANSPARENT)
    ax = fig.add_axes([0.1, 0.15, 0.85, 0.75])
    ax.set_facecolor(BOX_HEX)
    rates = [5, 8, 6, 10, 12, 35, 55, 70, 45, 30]
    colors = [ACCENT_HEX if r < 30 else RED_HEX for r in rates]
    ax.bar(np.arange(10), rates, color=colors, edgecolor=BOX_HEX, width=0.7)
    ax.axhline(y=30, color=ORANGE_HEX, ls="--", lw=2, label="30% threshold")
    ax.set_ylabel("Error Rate %", color=WHITE_HEX, fontsize=10)
    ax.set_xlabel("Time Window", color=WHITE_HEX, fontsize=10)
    ax.set_title("Anomaly Detection — Sliding Window", color=ACCENT_HEX, fontsize=12, fontweight="bold")
    ax.tick_params(colors=GRAY_HEX)
    ax.legend(facecolor=BOX_HEX, edgecolor=GRAY_HEX, labelcolor=WHITE_HEX)
    for spine in ax.spines.values(): spine.set_color(GRAY_HEX)
    return fig

def make_use_cases_diagram():
    fig, ax = plt.subplots(figsize=(7, 3.5), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 10); ax.set_ylim(0, 6); ax.axis("off")
    _box(ax, 3.5, 2.5, 3, 1, "Streaming\nMCP", ACCENT_HEX, 13)
    for x, y, label, c in [(0.2,4.5,"DevOps",ACCENT_HEX),(4,4.8,"Trading",ORANGE_HEX),(7.5,4.5,"Security",RED_HEX),
                            (0.2,0.5,"CI/CD",BLUE_HEX),(7.5,0.5,"IoT",ORANGE_HEX)]:
        _box(ax, x, y, 2.2, 0.8, label, c, 10)
        _arrow(ax, 5, 3 if y > 3 else 2.5, x+1.1, y+0.4+(-.4 if y > 3 else .4), c)
    return fig

def make_takeaways_diagram():
    fig, ax = plt.subplots(figsize=(7, 3), facecolor=TRANSPARENT)
    ax.set_facecolor("none"); ax.set_xlim(0, 10); ax.set_ylim(0, 5); ax.axis("off")
    for (x, y, label, sub), c in zip([(1,3.8,"Kafka","Durable\nStreaming"),(4,3.8,"WebSocket","Real-time\nPush"),
                                       (7,3.8,"Cedar","AuthZ\nPolicies"),(4,1,"Streaming\nMCP Server","")],
                                      [ORANGE_HEX, BLUE_HEX, RED_HEX, ACCENT_HEX]):
        _box(ax, x, y, 2, 1, f"{label}\n{sub}" if sub else label, c, 10)
    _arrow(ax, 2, 3.8, 4.2, 2.1, ORANGE_HEX); _arrow(ax, 5, 3.8, 5, 2.1, BLUE_HEX); _arrow(ax, 8, 3.8, 5.8, 2.1, RED_HEX)
    ax.text(5, 0.3, "Hybrid Push + Pull on a Single Connection", ha="center", color=ACCENT_HEX, fontsize=12, fontweight="bold")
    return fig


# ── Template layout indices ──────────────────────────────────────────────
# 0=TITLE, 1=SECTION_HEADER, 2=TITLE_AND_BODY, 3=1_Title and body,
# 4=TITLE_AND_TWO_COLUMNS, 5=TITLE_ONLY, 6=BLANK

# ── Slide definitions ────────────────────────────────────────────────────
# type: "title" -> layout 0 (PH0=title, PH1=subtitle)
# type: "section" -> layout 1 (PH0=title)
# type: "body" -> layout 2 (PH0=title, PH1=body)
# type: "two_col" -> layout 4 (PH0=title, PH1=left, PH2=right) — bullets left, diagram right
# type: "diagram_full" -> layout 5 (PH0=title) + image below

SLIDES = [
    {"type": "title", "title": "MCP Live:\nStreaming Context to AI Agents",
     "subtitle": "Building Reactive MCP Servers with Kafka, WebSockets & Cedar\ngithub.com/hakohli/mcp-streaming-demo",
     "notes": "TALKING POINTS:\n- Welcome everyone. Today we're extending MCP beyond request/response.\n- Core idea: what if MCP servers pushed real-time context to AI agents instead of waiting to be asked?\n- We'll walk through a working implementation using Kafka, WebSockets, and Cedar.\n- By the end, you'll see a live demo of an AI agent detecting a production incident in under 2 seconds.\n- Everything is open source and available in the repo."},

    {"type": "diagram_full", "title": "System Architecture",
     "diagram": make_full_architecture_diagram,
     "notes": "TALKING POINTS:\n- This is the full system architecture. Let me walk through each component.\n- Left: Log Simulator generates realistic logs from five services. Every ~30s it triggers an error spike.\n- Center-left: Apache Kafka (Confluent 7.6.0, KRaft mode) provides durability, replayability, and decoupling.\n- Center: The Streaming MCP Server has four sub-components: Kafka Consumer (rolling 100-event window), Anomaly Detection (error pattern tracking), WebSocket Server on port 8765 (pushes to subscribers with 500-event bounded queues), and Cedar Policy Engine (permit/forbid on every action, default deny).\n- Right: AI Agent subscribes once, runs a 20-event sliding window, flags anomalies at >30% error rate, and maps errors to fix suggestions.\n- Two arrow colors: teal = streaming push, blue = request/response pull. Both coexist on the same WebSocket."},

    {"type": "two_col", "title": "The Stale Context Problem",
     "bullets": ["Traditional MCP = request/response snapshots", "Agent asks for logs \u2192 gets point-in-time dump \u2192 stale data",
                 "Deploy at 2:03pm, errors at 2:04pm, agent polls at 2:15pm", "What if context came TO the agent, as it happened?"],
     "diagram": make_stale_context_diagram,
     "notes": "TALKING POINTS:\n- Today MCP works like a snapshot camera \u2014 the agent asks, gets a point-in-time answer, moves on.\n- Real scenario: deploy at 2:03pm, errors spike at 2:04pm, agent doesn't know until 2:15pm. That's an 11-minute blind spot.\n- The timeline diagram shows the red gap \u2014 the danger zone.\n- In incident response, trading, security \u2014 11 minutes is an eternity.\n- Can we flip the model so context flows TO the agent continuously?"},

    {"type": "two_col", "title": "Streaming MCP Architecture",
     "bullets": ["Standard MCP: Client \u2192 Server \u2192 Response (pull model)",
                 "Streaming MCP: Server \u2192 Client via subscriptions (push model)",
                 "Why Kafka: durable, replayable, backpressure handling",
                 "Hybrid protocol: streaming + request/response on same connection"],
     "diagram": make_architecture_diagram,
     "notes": "TALKING POINTS:\n- Top row: traditional pull-based MCP. Bottom row: what we're building.\n- Event sources feed Kafka, Kafka feeds the MCP server, server pushes to agents over WebSockets.\n- Why Kafka? (1) Durability \u2014 events persisted for replay. (2) Backpressure \u2014 handles slow consumers. (3) Already in your infra.\n- Key insight: hybrid protocol \u2014 streaming AND request/response on the same WebSocket. Zero changes to existing MCP tools."},

    {"type": "two_col", "title": "Key Protocol Extension: subscribe",
     "bullets": ["New method: subscribe / unsubscribe for continuous push",
                 "Still supports get_anomalies, get_context (request/response)",
                 "Bounded queues (500 events) handle backpressure automatically"],
     "diagram": make_subscribe_protocol_diagram,
     "notes": "TALKING POINTS:\n- Minimal extension \u2014 just two new methods: subscribe and unsubscribe.\n- Sequence: Agent sends subscribe('app-logs'). Events flow continuously. Error spike arrives immediately.\n- Mid-stream, agent sends get_anomalies() \u2014 hybrid protocol in action.\n- Backpressure: 500-event bounded queue per subscriber. Slow consumers get oldest events dropped.\n- Dead subscribers auto-cleaned \u2014 no resource leaks."},

    {"type": "two_col", "title": "Cedar Authorization for MCP",
     "bullets": ["Open-source policy language by AWS",
                 "Declarative: permit/forbid(principal, action, resource)",
                 "Role-based access control for MCP actions",
                 "Default deny \u2014 no matching policy = request denied"],
     "diagram": make_cedar_diagram,
     "notes": "TALKING POINTS:\n- Streaming data to agents raises the question: who's allowed to see what?\n- Cedar is open-source, now adopted in MCP ecosystem (Bedrock AgentCore, ToolHive).\n- Simple model: permit/forbid(principal, action, resource).\n- MCP actions: subscribe, get_anomalies, get_context, call_tool. Resources: streams like 'app-logs'.\n- Default deny \u2014 no matching policy = blocked. Secure by default."},

    {"type": "body", "title": "Cedar Policy Examples",
     "bullets": ['permit(principal, action == Action::"subscribe", resource == Stream::"app-logs");',
                 'permit(principal, action == Action::"get_anomalies", resource)\n    when { principal.role == "ops-agent" };',
                 'forbid(principal, action == Action::"subscribe", resource == Stream::"audit-logs")\n    when { principal.role == "readonly" };',
                 "Forbid always wins over permit \u2014 secure by default"],
     "font_size": 14,
     "notes": "TALKING POINTS:\n- Real Cedar policies from our demo.\n- First: anyone can subscribe to app-logs (general observability).\n- Second: only ops-agents can call get_anomalies (expensive operation).\n- Third: readonly agents explicitly forbidden from audit-logs (compliance data).\n- Forbid ALWAYS wins over permit. Human-readable, auditable, version-controllable."},

    {"type": "two_col", "title": "Live Demo: Real-time Log Monitoring",
     "bullets": ["Log simulator \u2192 Kafka \u2192 MCP Server \u2192 AI Agent",
                 "Error spike \u2192 agent detects anomaly within 1 second",
                 "Agent queries for anomaly summary mid-stream",
                 "Cedar blocks unauthorized agents from sensitive streams"],
     "diagram": make_demo_flow_diagram,
     "notes": "TALKING POINTS:\n- Stack: Confluent Kafka 7.6.0 in Docker, Python 3.9, confluent-kafka 2.13.2, websockets 15.0.1.\n- Log simulator generates realistic logs, pushes to Kafka. MCP server consumes and pushes to agents.\n- Normal traffic flows first. Then error spike hits.\n- Agent detected 'ANOMALY DETECTED \u2014 error rate 60% in last 20 events' within 1 second.\n- Automatic fixes: ConnectionRefused \u2192 check service; OOM \u2192 increase heap; Timeout \u2192 circuit breaker.\n- End-to-end: error to fix suggestion under 2 seconds."},

    {"type": "two_col", "title": "Anomaly Detection Flow",
     "bullets": ["Rolling window of last 100 events maintained server-side",
                 "Error rate > 30% triggers anomaly alert",
                 "Agent maps error patterns to fix suggestions automatically"],
     "diagram": make_anomaly_detection_diagram,
     "notes": "TALKING POINTS:\n- Server-side: rolling 100-event window for stable baseline.\n- Agent-side: sliding 20-event window, 30% threshold triggers alert.\n- Bar chart: green = normal, red = above threshold, orange line = trigger.\n- Agent calls get_anomalies() for structured summary with top-N error patterns.\n- Error \u2192 remediation mapping: ConnectionRefused \u2192 check connectivity; OOM \u2192 increase heap; SSL \u2192 renew certs.\n- All within 1-2 seconds. Compare to the 11-minute blind spot."},

    {"type": "body", "title": "Production Patterns & Gotchas",
     "bullets": ["1. Backpressure: bounded queues, drop slow subscribers",
                 "2. Hybrid protocol: streaming AND on-demand on same WebSocket",
                 "3. Kafka consumer groups: same = load balance, different = fan-out",
                 "4. Reconnection: exponential backoff + context snapshot",
                 "5. Be selective: stream errors & anomalies, batch normal metrics"],
     "notes": "TALKING POINTS:\n- Backpressure: 500-event asyncio queue per subscriber. Full = drop oldest. Consider dead-letter topics.\n- Hybrid protocol: don't force everything into one model. Subscribe for live feeds, query for summaries.\n- Consumer groups: same group ID = load balance (one consumer per event). Different = fan-out (all get every event).\n- Reconnection: exponential backoff to avoid thundering herd. Context snapshot on reconnect.\n- Be selective: stream errors and anomalies. Batch normal INFO metrics. Don't flood the pipe."},

    {"type": "two_col", "title": "What Streaming MCP Unlocks",
     "bullets": ["DevOps \u2014 real-time incident detection and auto-remediation",
                 "Trading \u2014 live market data feeding AI decision agents",
                 "Security \u2014 streaming audit logs with Cedar-authorized access",
                 "CI/CD \u2014 file watchers notifying agents of code changes",
                 "IoT \u2014 sensor data streaming to predictive maintenance agents"],
     "diagram": make_use_cases_diagram,
     "notes": "TALKING POINTS:\n- DevOps: what we demoed. Real-time incident detection, automated triage.\n- Trading: live market data, latency matters, polling is not an option.\n- Security: audit logs with Cedar access control. SOC agent gets alerts; readonly gets summaries.\n- CI/CD: file watchers trigger agents on code changes. Auto-test, security check, update docs.\n- IoT: sensor data to predictive maintenance. Detect anomalies before equipment fails.\n- Common thread: any domain where waiting to be asked is too slow."},

    {"type": "two_col", "title": "Key Takeaways",
     "bullets": ["MCP doesn't have to be request/response only",
                 "Kafka (Confluent 7.6.0) + WebSocket = durable, scalable streaming bridge",
                 "Cedar provides declarative, auditable authorization for MCP",
                 "Hybrid protocol (push + pull) on single connection is the sweet spot",
                 "All open source \u2014 Python 3.9+, Docker, ~200 lines of code"],
     "diagram": make_takeaways_diagram,
     "notes": "TALKING POINTS:\n- MCP is extensible beyond request/response. Streaming is a natural extension.\n- Kafka + WebSocket: battle-tested, widely deployed, well-understood.\n- Cedar: declarative, auditable policies. Default deny = start secure, open up explicitly.\n- Hybrid push + pull on one connection is the key design decision.\n- The diagram shows Kafka + WebSocket + Cedar converging into the streaming MCP server."},

    {"type": "title", "title": "Thank You & Q&A",
     "subtitle": "github.com/hakohli/mcp-streaming-demo\nAll open source: Kafka \u00b7 Cedar \u00b7 WebSockets \u00b7 Python",
     "notes": "TALKING POINTS:\n- Complete working code on GitHub: github.com/hakohli/mcp-streaming-demo\n- Stack: Confluent Kafka 7.6.0, Cedar, WebSockets 15.0.1, Python 3.9+.\n- To run: docker compose up -d, then streaming_mcp_server.py, log_simulator.py, agent_client.py.\n- Happy to take questions on architecture, Cedar policies, Kafka integration.\n- Interested in contributing? SSE transport, multi-topic subscriptions, persistent agent state \u2014 let's talk."},
]

# ── Build slides from template ───────────────────────────────────────────
prs = Presentation(TEMPLATE)

# Remove the 8 sample slides from the template
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

LAYOUTS = {name: i for i, name in enumerate(
    ["TITLE","SECTION_HEADER","TITLE_AND_BODY","1_Title and body",
     "TITLE_AND_TWO_COLUMNS","TITLE_ONLY","BLANK"])}

for s in SLIDES:
    stype = s["type"]

    if stype == "title":
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["TITLE"]])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = s["title"]
            elif ph.placeholder_format.idx == 1 and s.get("subtitle"):
                ph.text = s["subtitle"]

    elif stype == "section":
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["SECTION_HEADER"]])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = s["title"]

    elif stype == "body":
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["TITLE_AND_BODY"]])
        fs = s.get("font_size", 16)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = s["title"]
            elif ph.placeholder_format.idx == 1:
                tf = ph.text_frame
                for i, item in enumerate(s["bullets"]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = item
                    p.font.size = Pt(fs)
                    p.space_after = Pt(6)

    elif stype == "two_col":
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["TITLE_AND_TWO_COLUMNS"]])
        fs = s.get("font_size", 14)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = s["title"]
            elif ph.placeholder_format.idx == 1:
                tf = ph.text_frame
                for i, item in enumerate(s["bullets"]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = item
                    p.font.size = Pt(fs)
                    p.space_after = Pt(6)
            elif ph.placeholder_format.idx == 2 and s.get("diagram"):
                # Clear placeholder text, add image on top
                ph.text = ""
        if s.get("diagram"):
            fig = s["diagram"]()
            # Place diagram in right column area: 5.3in from left, 1.3in from top, 4.4x3.6in
            slide.shapes.add_picture(fig_to_stream(fig), Inches(5.3), Inches(1.3), Inches(4.4), Inches(3.6))

    elif stype == "diagram_full":
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUTS["TITLE_ONLY"]])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = s["title"]
        if s.get("diagram"):
            fig = s["diagram"]()
            # Full width below title: 0.3in from left, 1.1in from top
            slide.shapes.add_picture(fig_to_stream(fig), Inches(0.3), Inches(1.1), Inches(9.4), Inches(4.3))

    # Add speaker notes
    if s.get("notes"):
        slide.notes_slide.notes_text_frame.text = s["notes"]

out = "/Users/hakohli/mcp-streaming-demo/MCP_Live_Streaming_Context.pptx"
prs.save(out)
print(f"\u2705 Saved with DevSummit template: {out}")
